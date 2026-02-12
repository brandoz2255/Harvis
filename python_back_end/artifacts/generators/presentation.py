"""
PowerPoint presentation generator using python-pptx
"""

import os
import uuid
import logging
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)


def generate_presentation(content: Dict[str, Any], output_dir: str) -> str:
    """
    Generate PPTX file from manifest content.

    Args:
        content: Presentation content dict with 'slides' key
        output_dir: Directory to save the file

    Returns:
        Full path to generated file

    Expected content format:
    {
        "slides": [
            {
                "layout": "title",  # title, title_content, two_content, blank, section_header
                "title": "Slide Title",
                "content": ["Bullet point 1", "Bullet point 2"],  # for title_content
                "notes": "Speaker notes for this slide"
            }
        ],
        "title": "Presentation Title",
        "author": "Harvis AI",
        "theme": "default"  # optional
    }
    """
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()

    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Set presentation properties
    core_props = prs.core_properties
    core_props.author = content.get("author", "Harvis AI")
    if content.get("title"):
        core_props.title = content["title"]

    slides_data = content.get("slides", [])

    if not slides_data:
        # Create a default title slide
        slides_data = [{
            "layout": "title",
            "title": content.get("title", "Untitled Presentation")
        }]

    for slide_data in slides_data:
        layout_name = slide_data.get("layout", "title_content")
        slide_layout = _get_slide_layout(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Add title if present
        title_text = slide_data.get("title")
        if title_text and slide.shapes.title:
            slide.shapes.title.text = title_text

        # Add content based on layout
        if layout_name == "title":
            # Title slide - add subtitle if content provided
            subtitle_content = slide_data.get("content", [])
            if subtitle_content:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                        if isinstance(subtitle_content, list):
                            shape.text = "\n".join(str(item) for item in subtitle_content)
                        else:
                            shape.text = str(subtitle_content)
                        break

        elif layout_name == "title_content":
            # Title and content slide
            content_items = slide_data.get("content", [])
            if content_items:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # Content placeholder
                        tf = shape.text_frame
                        for i, item in enumerate(content_items):
                            if i == 0:
                                tf.text = str(item)
                            else:
                                p = tf.add_paragraph()
                                p.text = str(item)
                                p.level = 0
                        break

        elif layout_name == "two_content":
            # Two column content
            content_items = slide_data.get("content", [])
            if content_items:
                # Split content between two columns
                mid = len(content_items) // 2
                left_content = content_items[:mid] if mid > 0 else content_items
                right_content = content_items[mid:] if mid > 0 else []

                placeholders = list(slide.placeholders)
                content_placeholders = [p for p in placeholders if p.placeholder_format.idx > 0]

                if len(content_placeholders) >= 1 and left_content:
                    tf = content_placeholders[0].text_frame
                    for i, item in enumerate(left_content):
                        if i == 0:
                            tf.text = str(item)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(item)

                if len(content_placeholders) >= 2 and right_content:
                    tf = content_placeholders[1].text_frame
                    for i, item in enumerate(right_content):
                        if i == 0:
                            tf.text = str(item)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(item)

        elif layout_name == "section_header":
            # Section header - just title, maybe with subtitle
            content_items = slide_data.get("content", [])
            if content_items:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        if isinstance(content_items, list):
                            shape.text = "\n".join(str(item) for item in content_items)
                        else:
                            shape.text = str(content_items)
                        break

        # Add speaker notes if provided
        notes = slide_data.get("notes")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # Generate unique filename
    filename = f"artifact_{uuid.uuid4().hex[:16]}.pptx"
    filepath = os.path.join(output_dir, filename)

    prs.save(filepath)
    logger.info(f"Generated presentation: {filepath}")

    return filepath


def _get_slide_layout(prs: Presentation, layout_name: str):
    """Get slide layout by name"""
    layout_map = {
        "title": 0,           # Title Slide
        "title_content": 1,   # Title and Content
        "section_header": 2,  # Section Header
        "two_content": 3,     # Two Content
        "comparison": 4,      # Comparison
        "title_only": 5,      # Title Only
        "blank": 6,           # Blank
        "content_caption": 7, # Content with Caption
        "picture_caption": 8, # Picture with Caption
    }

    layout_idx = layout_map.get(layout_name, 1)  # Default to title_content

    # Get layout from default presentation
    try:
        return prs.slide_layouts[layout_idx]
    except IndexError:
        # Fall back to first available layout
        return prs.slide_layouts[0]
